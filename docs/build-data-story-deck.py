#!/usr/bin/env python3
"""Store All Barbados — data-story deck in the dashboard's design theme.

Design language mirrors index.html:
  red #E30613 accent, charcoal #111/#222 dark UI, white cards with #e4e4e4
  hairlines, #f2f2f2 tracks, Montserrat display type, uppercase letterspaced
  eyebrows, red-circled section numbers, dark insight boxes with a red left rule.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------------- theme ----------------
ACCENT   = RGBColor(0xE3, 0x06, 0x13)
ACCENT_D = RGBColor(0xB0, 0x00, 0x10)   # --teal in the dashboard (dark red)
PINK     = RGBColor(0xF6, 0xA3, 0xA3)
PINKPALE = RGBColor(0xFF, 0xD9, 0xD9)
INK      = RGBColor(0x11, 0x11, 0x11)
CHAR     = RGBColor(0x22, 0x22, 0x22)
MUT      = RGBColor(0x6B, 0x6B, 0x6B)
LINE     = RGBColor(0xE4, 0xE4, 0xE4)
PAPER2   = RGBColor(0xF2, 0xF2, 0xF2)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GOOD     = RGBColor(0x2F, 0x7D, 0x5B)
WARN     = RGBColor(0xB8, 0x86, 0x20)
G_MID    = RGBColor(0x7A, 0x7A, 0x7A)
G_LT     = RGBColor(0xB5, 0xB5, 0xB5)
G_XLT    = RGBColor(0xD5, 0xD5, 0xD5)
DARKBG   = RGBColor(0x14, 0x14, 0x14)
DARKCARD = RGBColor(0x1F, 0x1F, 0x1F)
DARKLINE = RGBColor(0x3A, 0x3A, 0x3A)
FB       = RGBColor(0x3B, 0x59, 0x98)

HEAD = "Montserrat"
BODY = "Arial"          # metric-safe stand-in for Source Sans 3

EMU = 914400
SW, SH = 13.333, 7.5
ML = 0.55               # left/right margin

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]

# ---------------- data (aggregated from index.html LOCATIONS et al.) ----------------
# occupancy
OCC = dict(total=985, occupied=976, rentable=982, vacant_rentable=6, oos=14,
           rent=273729, gross=283461, mi=34, mo=13, net=21)
SITES_OCC = [  # name, occupied, rentable, vacant, note
    ("Central", 457, 461, 4,  "The Pine, St. Michael"),
    ("South",   264, 265, 1,  "Gibbons, Christ Church"),
    ("Lears",   255, 256, 1,  "14 more units out of service"),
]
FUNNEL = dict(inquiries=361, converted=138)
# payments — June receipts, all sites
PAY_TIERS = [  # label, units, amount, color
    ("Autopay (Plug & Pay)", 274,  89148, ACCENT),
    ("Online / bank app",    343, 149239, CHAR),
    ("Manual card",           82,  41040, G_MID),
    ("Cash / cheque",        113,  58990, G_LT),
]
PAY_TOT_U = sum(t[1] for t in PAY_TIERS)          # 812
PAY_TOT_A = sum(t[2] for t in PAY_TIERS)          # 338,417
PAY_METHODS = [  # label, units, amount
    ("Plug & Pay (autopay)",       274,  89148),
    ("CIBC bank-app transfer",     228, 105615),
    ("Website Pay Now",            110,  41415),
    ("Cash",                        71,  30441),
    ("Visa (manual)",               55,  32104),
    ("Cheque",                      42,  28550),
    ("Mastercard (manual)",         27,   8936),
    ("Republic bank-app transfer",   5,   2210),
]
AP_TREND = [("Jan", 27.4), ("Feb", 27.7), ("Mar", 28.3),
            ("Apr", 28.6), ("May", 28.5), ("Jun", 28.0)]
# acquisition — marketing source (974 tenants) and lifetime value
SRC = [  # label, tenants, value, color
    ("Internet",        326,  3141552, ACCENT),
    ("Referral",        217,  3176225, CHAR),
    ("Previous tenant", 160,  1997545, ACCENT_D),
    ("Drive-by",         84,  1006195, G_MID),
    ("Yellow Pages",     60,  1270761, WARN),
    ("Facebook",          6,    65338, FB),
    ("Other",           103,  1322899, G_LT),
    ("Unknown",          18,   550800, G_XLT),
]
SRC_TOT   = sum(s[1] for s in SRC)                 # 974
VAL_TOT   = sum(s[2] for s in SRC)                 # 12,531,315
# contact channel (979 units)
CHAN = [
    ("Phone",        602, CHAR),
    ("Not recorded", 165, G_XLT),
    ("Walk-in",      148, G_MID),
    ("Email",         47, ACCENT_D),
    ("Web",           17, ACCENT),
]
CHAN_TOT = sum(c[1] for c in CHAN)                 # 979
# customer profile (978 active)
WHY    = [("Location", 657), ("Features", 179), ("Management", 77), ("Other", 44), ("Price", 2)]
REASON = [("Excess / declutter", 338), ("Moving", 249), ("Business needs", 231),
          ("Renovating", 53), ("Other", 88)]

# ---------------- helpers ----------------
def _set_spc(run, hundredths):
    run.font._rPr.set('spc', str(hundredths))

def _strip_style(s):
    """remove <p:style> so the theme's effectRef shadow never applies."""
    el = s._element
    st = el.find(qn('p:style'))
    if st is not None:
        el.remove(st)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, round_=None, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_ is not None:
        try: s.adjustments[0] = round_
        except Exception: pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    _strip_style(s)
    return s

def text(slide, x, y, w, h, runs, size=12, color=INK, font=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line_spacing=None,
         wrap=True, shrink=False):
    """runs: str, or list of (txt, dict-overrides) tuples; \n splits paragraphs."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    paras, cur = [], []
    for txt, ov in runs:
        parts = txt.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                paras.append(cur); cur = []
            if part:
                cur.append((part, ov))
    paras.append(cur)
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for txt, ov in para:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = ov.get('font', font)
            f.size = Pt(ov.get('size', size))
            f.bold = ov.get('bold', bold)
            f.color.rgb = ov.get('color', color)
            if ov.get('italic'): f.italic = True
            sp = ov.get('spc', spc)
            if sp: _set_spc(r, sp)
    return tb

def chip(slide, x, y, num, dark=False):
    """red-outlined circular section number (dashboard .sec-num)."""
    d = 0.42
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = DARKCARD if dark else WHITE
    c.line.color.rgb = ACCENT; c.line.width = Pt(1.5)
    c.shadow.inherit = False
    _strip_style(c)
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = HEAD; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = ACCENT

def header(slide, num, eyebrow, title, sub=None):
    chip(slide, ML, 0.42, num)
    text(slide, ML + 0.58, 0.40, 9.0, 0.3, eyebrow.upper(), size=11, font=HEAD,
         bold=True, color=ACCENT, spc=220)
    text(slide, ML + 0.58, 0.63, 11.6, 0.62, title, size=27, font=HEAD, bold=True, color=INK)
    y = 1.22
    if sub:
        text(slide, ML + 0.58, y, 11.4, 0.3, sub, size=11.5, color=MUT)
        y += 0.34
    rect(slide, ML, y + 0.04, SW - 2 * ML, 0.012, fill=LINE)
    return y + 0.18

def footer(slide, n, note, dark=False):
    col = RGBColor(0x8A, 0x8A, 0x8A) if dark else MUT
    text(slide, ML, 7.08, 10.6, 0.25,
         [("Store All Barbados — Data Story · July 2026", {'bold': True}),
          ("   ·   " + note, {})], size=8.5, color=col)
    text(slide, SW - ML - 0.6, 7.08, 0.6, 0.25, f"{n:02d}", size=9, font=HEAD,
         bold=True, color=col, align=PP_ALIGN.RIGHT)

def stat_tile(slide, x, y, w, h, k, v, v_small, d, accent=False, dark=False):
    """dashboard .stat card: uppercase key, huge Montserrat value, caption."""
    rect(slide, x, y, w, h,
         fill=DARKCARD if dark else WHITE,
         line=DARKLINE if dark else LINE, line_w=1, round_=0.09)
    pad = 0.17
    text(slide, x + pad, y + 0.14, w - 2 * pad, 0.24, k.upper(), size=9, font=HEAD,
         bold=True, color=RGBColor(0xA8, 0xA8, 0xA8) if dark else MUT, spc=110)
    vcol = ACCENT if accent else (WHITE if dark else INK)
    text(slide, x + pad, y + 0.40, w - 2 * pad, 0.55,
         [(v, {}), (v_small, {'size': 14, 'color': PINK if dark else MUT, 'bold': False})],
         size=30, font=HEAD, bold=True, color=vcol)
    text(slide, x + pad, y + 0.97, w - 2 * pad, h - 1.02, d, size=9.5,
         color=RGBColor(0xC0, 0xC0, 0xC0) if dark else MUT, line_spacing=1.05)

def insight(slide, x, y, w, h, title, body_runs, bg=CHAR):
    """dashboard .insight: dark box, 4px red left rule, pale-pink heading."""
    rect(slide, x, y, w, h, fill=bg, round_=0.045)
    rect(slide, x, y + 0.05, 0.055, h - 0.10, fill=ACCENT)
    pad = 0.26
    text(slide, x + pad, y + 0.20, w - 2 * pad, 0.34, title, size=15, font=HEAD,
         bold=True, color=PINKPALE)
    text(slide, x + pad, y + 0.58, w - 2 * pad, h - 0.74, body_runs, size=11,
         color=RGBColor(0xE0, 0xE0, 0xE0), line_spacing=1.18)

def card(slide, x, y, w, h, title=None):
    rect(slide, x, y, w, h, fill=WHITE, line=LINE, line_w=1, round_=0.05)
    if title:
        rect(slide, x + 0.22, y + 0.235, 0.09, 0.09, fill=ACCENT, round_=0.5)
        text(slide, x + 0.38, y + 0.175, w - 0.6, 0.26, title.upper(), size=9.5,
             font=HEAD, bold=True, color=CHAR, spc=90)

def bar_row(slide, x, y, w, label, frac, val, color=ACCENT, lab_w=1.55, val_w=0.95,
            bar_h=0.125, lab_size=10, muted_val=False):
    """dashboard .bar-row: label · track+fill · right-aligned number."""
    track_w = w - lab_w - val_w - 0.22
    text(slide, x, y - 0.035, lab_w, 0.24, label, size=lab_size, bold=True, color=INK)
    rect(slide, x + lab_w + 0.08, y, track_w, bar_h, fill=PAPER2, round_=0.5)
    fw = max(track_w * min(frac, 1.0), 0.055)
    rect(slide, x + lab_w + 0.08, y, fw, bar_h, fill=color, round_=0.5)
    text(slide, x + w - val_w, y - 0.035, val_w, 0.24, val, size=9.5, bold=True,
         color=MUT if muted_val else CHAR, align=PP_ALIGN.RIGHT)

def stacked_bar(slide, x, y, w, h, parts, total, min_label_frac=0.075):
    """single horizontal 100% bar with 2px white gaps; % inside wide segments."""
    gap = 0.028
    cx = x
    for label, n, color in parts:
        seg = w * n / total - gap
        if seg <= 0.01:
            seg = 0.012
        rect(slide, cx, y, seg, h, fill=color, round_=0.18)
        if n / total >= min_label_frac:
            lum = 0.2126 * color[0] + 0.7152 * color[1] + 0.0722 * color[2]
            tcol = WHITE if lum < 140 else CHAR
            text(slide, cx, y + h / 2 - 0.115, seg, 0.23, f"{100*n/total:.0f}%",
                 size=10.5, font=HEAD, bold=True, color=tcol, align=PP_ALIGN.CENTER, wrap=False)
        cx += seg + gap

def legend_row(slide, x, y, w, items, per_row=4, size=9.5):
    cw = w / per_row
    for i, (label, sub, color) in enumerate(items):
        ix = x + (i % per_row) * cw
        iy = y + (i // per_row) * 0.26
        rect(slide, ix, iy + 0.035, 0.13, 0.13, fill=color, round_=0.25)
        text(slide, ix + 0.20, iy - 0.015, cw - 0.24, 0.24,
             [(label, {'bold': True, 'color': INK}),
              ("  " + sub, {'color': MUT, 'size': size - 0.5})], size=size, wrap=False)

def dark_bg(slide):
    rect(slide, 0, 0, SW, SH, fill=DARKBG)
    rect(slide, 0, 0, SW, 0.07, fill=ACCENT)            # header::before red rule
    rect(slide, 8.4, 0, 4.94, 0.07, fill=ACCENT_D)

def fmt_k(n): return f"${n/1000:,.0f}k"

# ================= SLIDE 1 — TITLE =================
s = prs.slides.add_slide(BLANK)
dark_bg(s)
text(s, ML, 0.9, 12, 0.3, "STORE ALL BARBADOS  ·  PERFORMANCE REPORTING", size=12,
     font=HEAD, bold=True, color=PINK, spc=300)
text(s, ML, 1.28, 12.3, 1.9,
     [("Full house. Digital money.\n", {}),
      ("Word of mouth", {'color': ACCENT}),
      (" pays the rent.", {})],
     size=42, font=HEAD, bold=True, color=WHITE, line_spacing=1.04)
text(s, ML, 3.05, 9.6, 0.6,
     "What the June 2026 data says about where Store All grows next — told in four acts, "
     "from the dashboard: occupancy, payments, acquisition and the customer.",
     size=13.5, color=RGBColor(0xD8, 0xD8, 0xD8), line_spacing=1.25)
acts = [
    ("01", "OCCUPANCY", "We are sold out — 6 units left across 985."),
    ("02", "PAYMENTS", "86% of payers are already digital; autopay is stuck."),
    ("03", "ACQUISITION", "Referrals carry the value; the phone carries the volume."),
    ("04", "CUSTOMER", "Local, residential, location-led — never price-led."),
]
aw = (SW - 2 * ML - 3 * 0.22) / 4
for i, (n, t, d) in enumerate(acts):
    ax = ML + i * (aw + 0.22)
    rect(s, ax, 4.35, aw, 1.72, fill=DARKCARD, line=DARKLINE, line_w=1, round_=0.07)
    chip(s, ax + 0.2, 4.55, n, dark=True)
    text(s, ax + 0.2, 5.08, aw - 0.4, 0.26, t, size=11.5, font=HEAD, bold=True,
         color=WHITE, spc=120)
    text(s, ax + 0.2, 5.38, aw - 0.4, 0.62, d, size=9.5,
         color=RGBColor(0xB8, 0xB8, 0xB8), line_spacing=1.12)
rect(s, ML, 6.45, SW - 2 * ML, 0.012, fill=DARKLINE)
text(s, ML, 6.58, 12, 0.3,
     "Central · South · Lears   —   SiteLink rent roll, receipts, inquiry tracking & marketing summary · June 2026 month-end",
     size=10, color=RGBColor(0x9A, 0x9A, 0x9A))
footer(s, 1, "Title", dark=True)

# ================= SLIDE 2 — ACT 01a: SOLD OUT =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "01", "Act 01 · Occupancy", "We are sold out.",
            "June 2026 month-end · 976 of 982 rentable units occupied across three sites (99.4%)")
tiles = [
    ("Rentable occupancy", "99.4", "%", "976 of 982 rentable units filled — physical occupancy 99.1% of all 985", True),
    ("Units left to sell", "6", "", "vacant rentable units in the whole portfolio: 4 Central · 1 South · 1 Lears", False),
    ("Monthly rent roll", "$274", "k", "96.6% of the $283k gross potential — near-zero discounting", False),
    ("Net absorption · June", "+21", "", "34 move-ins vs 13 move-outs — vacated units refill immediately", False),
]
tw = (SW - 2 * ML - 3 * 0.2) / 4
for i, (k, v, vs, d, acc) in enumerate(tiles):
    stat_tile(s, ML + i * (tw + 0.2), y0 + 0.06, tw, 1.62, k, v, vs, d, accent=acc)
cy = y0 + 1.88
card(s, ML, cy, 7.35, 2.55, "Occupancy by site — share of rentable units")
by = cy + 0.62
for name, occ, rentable, vac, note in SITES_OCC:
    frac = occ / rentable
    bar_row(s, ML + 0.25, by, 6.85, name, frac, f"{100*frac:.1f}%", color=ACCENT,
            lab_w=1.05, val_w=0.8, bar_h=0.16, lab_size=11)
    text(s, ML + 1.38, by + 0.20, 5.0, 0.2,
         f"{occ} of {rentable} rentable · {vac} vacant — {note}", size=8.5, color=MUT)
    by += 0.60
insight(s, ML + 7.55, cy, SW - ML - (ML + 7.55), 2.55, "Churn at the ceiling",
        [("A small negative month is tenants swapping, not demand falling: every unit "
          "vacated in June was re-let within the month. Economic occupancy of ", {}),
         ("96.6%", {'bold': True, 'color': WHITE}),
         (" means we collect close to full list price — we are not buying this occupancy "
          "with discounts.\n", {}),
         ("The portfolio has been at this ceiling all year. That is the story of this "
          "deck: ", {}),
         ("the constraint is no longer demand — it is space.", {'bold': True, 'color': WHITE})])
footer(s, 2, "SiteLink rent roll · June month-end; Lears position as of 2 July 2026")

# ================= SLIDE 3 — ACT 01b: WE NEED SUPPLY =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "01", "Act 01 · Occupancy", "Demand keeps arriving. Supply doesn't.",
            "SiteLink inquiry tracking · leads placed January–June 2026, all three sites")
cy = y0 + 0.06
card(s, ML, cy, 7.35, 2.9, "Six months of demand vs what we can sell today")
iy = cy + 0.62
rows = [
    ("Inquiries · Jan–Jun", FUNNEL['inquiries'], CHAR),
    ("Converted to move-ins", FUNNEL['converted'], ACCENT_D),
    ("Rentable units vacant today", OCC['vacant_rentable'], ACCENT),
]
mx = FUNNEL['inquiries']
for label, n, colr in rows:
    text(s, ML + 0.25, iy - 0.03, 2.35, 0.26, label, size=10.5, bold=True, color=INK)
    track_w = 3.6
    rect(s, ML + 2.7, iy, track_w, 0.24, fill=PAPER2, round_=0.5)
    rect(s, ML + 2.7, iy, max(track_w * n / mx, 0.07), 0.24, fill=colr, round_=0.5)
    text(s, ML + 6.4, iy - 0.05, 0.7, 0.3, str(n), size=15, font=HEAD, bold=True,
         color=colr if colr == ACCENT else CHAR, align=PP_ALIGN.RIGHT)
    iy += 0.52
text(s, ML + 0.25, iy + 0.06, 6.9, 0.75,
     [("361 inquiries chased 6 available units.", {'bold': True, 'color': INK}),
      (" Lears — the newest site — absorbed +15 net units in June alone and is now "
       "effectively full. The 38% inquiry→move-in rate is demand we largely turn away.", {})],
     size=10.5, color=MUT, line_spacing=1.2)
insight(s, ML + 7.55, cy, SW - ML - (ML + 7.55), 2.9, "We need supply.",
        [("Growth is capped by capacity, not by the market. At 99.4% full, revenue can "
          "only move through ", {}),
         ("rate", {'bold': True, 'color': WHITE}),
         (" — volume is gone until we add space.\n", {}),
         ("Every act that follows — payments, referrals, the website — makes the "
          "machine more efficient. But only ", {}),
         ("new units", {'bold': True, 'color': WHITE}),
         (" let it grow. Expansion is decision #1.", {})])
ky = cy + 3.14
rect(s, ML, ky, SW - 2 * ML, 1.02, fill=ACCENT, round_=0.09)
text(s, ML + 0.35, ky + 0.17, 11.6, 0.42,
     [("The ask:   ", {'color': PINKPALE, 'size': 13}),
      ("commit to a supply decision — expand Lears or open site #4.", {})],
     size=18, font=HEAD, bold=True, color=WHITE, wrap=False)
text(s, ML + 0.35, ky + 0.62, 11.6, 0.28,
     "Where the catchment data says it will fill is Act 04.",
     size=10, color=PINKPALE)
footer(s, 3, "SiteLink inquiry tracking Jan–Jun 2026 · rent roll June month-end")

# ================= SLIDE 4 — ACT 02a: MONEY IS DIGITAL =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "02", "Act 02 · Payments", "The money is already digital.",
            "June 2026 receipts, all sites · bank-app transfers reclassified July 2026 — they were never cash")
tiles = [
    ("Payers who pay electronically", "86", "%", "699 of 812 units paying in June — autopay, bank app, website or card", True),
    ("Revenue collected digitally", "82.6", "%", f"{fmt_k(PAY_TOT_A-58990)} of {fmt_k(PAY_TOT_A)} June receipts", False),
    ("Biggest single money channel", "$106", "k", "CIBC bank-app transfers — 31% of collections, ahead of autopay", False),
    ("Still paying cash", "71", " units", "just 9% of collections ($30k) — cash is the tail, not the norm", False),
]
tw = (SW - 2 * ML - 3 * 0.2) / 4
for i, (k, v, vs, d, acc) in enumerate(tiles):
    stat_tile(s, ML + i * (tw + 0.2), y0 + 0.06, tw, 1.62, k, v, vs, d, accent=acc)
cy = y0 + 1.88
card(s, ML, cy, SW - 2 * ML, 2.55, "How the 812 June payers paid — units by tier")
stacked_bar(s, ML + 0.25, cy + 0.62, SW - 2 * ML - 0.5, 0.5,
            [(l, n, c) for l, n, a, c in PAY_TIERS], PAY_TOT_U)
legend_row(s, ML + 0.25, cy + 1.28, SW - 2 * ML - 0.5,
           [(l, f"{n} units · {fmt_k(a)}/mo", c) for l, n, a, c in PAY_TIERS], per_row=4)
text(s, ML + 0.25, cy + 1.72, SW - 2 * ML - 0.5, 0.66,
     [("What changed: ", {'bold': True, 'color': INK}),
      ("July's reclassification revealed that “Money Order” and “Internet” receipts were CIBC and "
       "Republic bank-app transfers, and “Bank Transfer” was the website's Pay Now. The cash-heavy picture "
       "was a labelling illusion — customers moved to digital years ago. The gap left is not ", {}),
      ("digital vs cash", {'italic': True}),
      (" — it is ", {}),
      ("automatic vs manual", {'bold': True, 'color': ACCENT}),
      (", which is the next slide.", {})],
     size=10.5, color=MUT, line_spacing=1.22)
footer(s, 4, "SiteLink receipts June 2026 · 812 of 976 occupied units transacted in the window")

# ================= SLIDE 5 — ACT 02b: AUTOPAY STUCK =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "02", "Act 02 · Payments", "Autopay is stuck. 425 units are one click away.",
            "Plug & Pay share of 978 occupied units, January–June 2026 · June receipts by conversion tier")
cy = y0 + 0.06
card(s, ML, cy, 5.6, 3.0, "Autopay adoption — six flat months")
base_y = cy + 2.28; max_h = 1.45; ceil = 32.0
bw = (5.6 - 0.5 - 5 * 0.14) / 6
for i, (m, p) in enumerate(AP_TREND):
    bx = ML + 0.25 + i * (bw + 0.14)
    bh = max_h * p / ceil
    rect(s, bx, base_y - bh, bw, bh, fill=ACCENT if m == "Jun" else CHAR, round_=0.12)
    text(s, bx - 0.1, base_y - bh - 0.24, bw + 0.2, 0.2, f"{p:.1f}%", size=9,
         font=HEAD, bold=True, color=CHAR, align=PP_ALIGN.CENTER, wrap=False)
    text(s, bx - 0.1, base_y + 0.05, bw + 0.2, 0.2, m, size=9, color=MUT,
         align=PP_ALIGN.CENTER, wrap=False)
text(s, ML + 0.25, base_y + 0.32, 5.1, 0.3,
     [("+0.6pp in six months.", {'bold': True, 'color': INK}),
      (" Organic growth will not shift the mix — it needs a push.", {})],
     size=10, color=MUT)
card(s, ML + 5.8, cy, 7.0 - 0.02, 3.0, "The conversion ladder — June payers by distance from autopay")
ly = cy + 0.60
ladder = [
    ("On autopay today", 274, "Plug & Pay — the destination", ACCENT),
    ("One click away", 425, "343 bank-app / Pay Now + 82 manual card — already digital, just manual", CHAR),
    ("The hard tail", 113, "cash & cheque — needs an ACH offer at renewal, not a nudge", G_LT),
]
for label, n, d, colr in ladder:
    rect(s, ML + 6.05, ly, 0.09, 0.52, fill=colr)
    text(s, ML + 6.28, ly - 0.02, 3.3, 0.28, label, size=11.5, font=HEAD, bold=True, color=INK)
    text(s, ML + 6.28, ly + 0.26, 5.35, 0.26, d, size=9, color=MUT)
    text(s, ML + 11.35, ly - 0.04, 1.15, 0.4, str(n), size=21, font=HEAD, bold=True,
         color=colr if colr != G_LT else MUT, align=PP_ALIGN.RIGHT)
    ly += 0.66
text(s, ML + 6.05, ly + 0.02, 6.5, 0.3,
     [("Convert the one-click tier and autopay reaches ", {}),
      ("~72% of occupied units", {'bold': True, 'color': ACCENT}),
      (".", {})], size=10.5, color=MUT)
iy = cy + 3.22
insight(s, ML, iy, SW - 2 * ML, 1.12, "The ask: run the switch campaign",
        [("One targeted message to the 343 bank-app & Pay Now payers (“make it recurring”), a tick-box prompt "
          "for the 82 manual card payers at their next payment, and an ACH offer to cash & cheque at renewal. "
          "Every converted unit removes a monthly manual touch and cuts late-payment risk.", {})])
footer(s, 5, "SiteLink receipts Jan–Jun 2026 · autopay share of occupied units")

# ================= SLIDE 6 — ACT 03a: WORD OF MOUTH =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "03", "Act 03 · Acquisition — how they found us", "Word of mouth is the value engine.",
            "SiteLink marketing category · 974 occupied tenants · lifetime value = cumulative rent per tenant, $12.5M total")
cy = y0 + 0.06
card(s, ML, cy, SW - 2 * ML, 1.62, "How the current base found Store All — 974 tenants")
stacked_bar(s, ML + 0.25, cy + 0.56, SW - 2 * ML - 0.5, 0.42,
            [(l, n, c) for l, n, v, c in SRC], SRC_TOT, min_label_frac=0.06)
SHORT = {"Previous tenant": "Prev tenant", "Yellow Pages": "Y. Pages"}
legend_row(s, ML + 0.25, cy + 1.10, SW - 2 * ML - 0.5,
           [(SHORT.get(l, l), f"{n} · {100*n/SRC_TOT:.0f}%", c) for l, n, v, c in SRC],
           per_row=8, size=8)
c2 = cy + 1.84
card(s, ML, c2, 7.35, 2.6, "Lifetime value per tenant, by source")
vy = c2 + 0.56
vrows = [("Referral", 217, 14637, ACCENT), ("Previous tenant", 160, 12485, ACCENT_D),
         ("Drive-by", 84, 11979, G_MID), ("Internet", 326, 9637, CHAR)]
for label, n, per, colr in vrows:
    bar_row(s, ML + 0.25, vy, 6.85, label, per / 15000, f"${per:,.0f}", color=colr,
            lab_w=1.6, val_w=0.95, bar_h=0.15, lab_size=10.5)
    vy += 0.38
text(s, ML + 0.25, vy + 0.02, 6.85, 0.42,
     "Yellow Pages tops the table at ~$21k/tenant but is a shrinking legacy book of 60 "
     "long-tenured commercial accounts. Internet brings the most tenants — at the lowest value per head.",
     size=8.5, color=MUT, line_spacing=1.15)
insight(s, ML + 7.55, c2, SW - ML - (ML + 7.55), 2.6, "Referrals punch above their weight",
        [("Referral + previous tenant are ", {}),
         ("39% of tenants but 41% of all lifetime value — $5.2M", {'bold': True, 'color': WHITE}),
         (". A referred tenant is worth ", {}),
         ("$14.6k", {'bold': True, 'color': WHITE}),
         (", 52% more than an internet-sourced one.\n", {}),
         ("And we currently pay nothing for it. ", {}),
         ("The ask:", {'bold': True, 'color': PINKPALE}),
         (" formalise it — a referral credit for existing tenants and a win-back offer at move-out.", {})])
footer(s, 6, "SiteLink marketing summary (sMktCat) · occupied base + cumulative rent value")

# ================= SLIDE 7 — ACT 03b: PHONE / WEBSITE =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "03", "Act 03 · Acquisition — how they reached us", "They find us online — then they phone.",
            "SiteLink enquiry channel · 979 occupied units · contact channel, distinct from the marketing source")
tiles = [
    ("Contacted us by phone", "61", "%", "602 of 979 tenants — the phone is the front door at every site", False),
    ("Found us on the internet", "34", "%", "332 tenants incl. Facebook — discovery is already digital", False),
    ("Ever booked via the website", "1.7", "%", "17 tenants in the entire base — most predate 2020", True),
    ("Web bookings, last 12 months", "1", " of 270", "one web move-in out of 270 recent — the funnel is broken", True),
]
tw = (SW - 2 * ML - 3 * 0.2) / 4
for i, (k, v, vs, d, acc) in enumerate(tiles):
    stat_tile(s, ML + i * (tw + 0.2), y0 + 0.06, tw, 1.62, k, v, vs, d, accent=acc)
cy = y0 + 1.88
card(s, ML, cy, 7.35, 2.55, "Contact channel mix — 979 occupied units")
by = cy + 0.62
for label, n, colr in CHAN:
    bar_row(s, ML + 0.25, by, 6.85, label, n / CHAN_TOT / 0.65,
            f"{n} · {100*n/CHAN_TOT:.0f}%", color=colr, lab_w=1.5, val_w=1.1,
            bar_h=0.15, lab_size=10.5)
    by += 0.375
insight(s, ML + 7.55, cy, SW - ML - (ML + 7.55), 2.55, "A brochure, not a funnel",
        [("A third of tenants ", {}),
         ("discover", {'italic': True}),
         (" us online — then pick up the phone to book. Discovery is digital; booking is analog. "
          "That gap is the website.\n", {}),
         ("The ask:", {'bold': True, 'color': PINKPALE}),
         (" rebuild it as a booking channel — live availability, prices, “Reserve now” — and own local "
          "search (Google Business Profile, click-to-call). Its Pay Now already collects "
          "$41k/mo; the front of the site should work as hard as the back.", {})])
footer(s, 7, "SiteLink InquirySource · occupied base; recent = move-ins in the last 12 months")

# ================= SLIDE 8 — ACT 04: CUSTOMER PROFILE =================
s = prs.slides.add_slide(BLANK)
y0 = header(s, "04", "Act 04 · Customer profile", "Here for location — never for price.",
            "Local households in transition · SiteLink marketing summary · 978 active tenants, all sites")
tiles = [
    ("Chose us for location", "67", "%", "657 tenants — versus 2 (0.2%) for price. Convenience, not cost", True),
    ("Residential customers", "76", "%", "747 households vs 231 commercial accounts", False),
    ("Store furniture / household", "64", "%", "629 tenants — life-event storage dominates the mix", False),
    ("Live within ~5 miles", "43", "%", "and 13% overseas — the diaspora, concentrated at Central", False),
]
tw = (SW - 2 * ML - 3 * 0.2) / 4
for i, (k, v, vs, d, acc) in enumerate(tiles):
    stat_tile(s, ML + i * (tw + 0.2), y0 + 0.06, tw, 1.62, k, v, vs, d, accent=acc)
cy = y0 + 1.88
card(s, ML, cy, 5.95, 2.55, "Why they store")
by = cy + 0.60
r_tot = sum(n for _, n in REASON)
for label, n in REASON:
    bar_row(s, ML + 0.25, by, 5.45, label, n / 338, f"{n} · {100*n/r_tot:.0f}%",
            color=CHAR, lab_w=1.55, val_w=1.0, bar_h=0.14, lab_size=10)
    by += 0.375
card(s, ML + 6.15, cy, 3.1, 2.55, "Why they chose us")
by = cy + 0.60
w_tot = sum(n for _, n in WHY)
for label, n in WHY:
    bar_row(s, ML + 6.38, by, 2.68, label, n / 657, f"{100*n/w_tot:.0f}%",
            color=ACCENT, lab_w=1.05, val_w=0.42, bar_h=0.13, lab_size=9)
    by += 0.375
insight(s, ML + 9.45, cy, SW - ML - (ML + 9.45), 2.55, "Back to Act 01",
        [("Demand is local, residential and location-led — people store because life "
          "changed, near where they live.\n", {}),
         ("So the growth lever is ", {}),
         ("proximity and visibility", {'bold': True, 'color': WHITE}),
         (", not discounting — and it is exactly why ", {}),
         ("new supply in the right catchment", {'bold': True, 'color': WHITE}),
         (" will fill.", {})])
footer(s, 8, "SiteLink marketing summary · shares of tenants with the field recorded")

# ================= SLIDE 9 — CLOSING =================
s = prs.slides.add_slide(BLANK)
dark_bg(s)
text(s, ML, 0.55, 12, 0.3, "THE STORY, IN ONE SLIDE", size=11, font=HEAD, bold=True,
     color=PINK, spc=260)
text(s, ML, 0.85, 12.3, 0.65, "What the data asks of us", size=30, font=HEAD,
     bold=True, color=WHITE)
asks = [
    ("01", "Secure new supply", "Sold out at 99.4% with 361 inquiries chasing 6 units. "
     "Volume growth is gone until we add space — expand Lears or commit to site #4.",
     "EXPANSION"),
    ("02", "Convert the one-click 425", "Autopay is stuck at 28% (+0.6pp in 6 months). "
     "425 units already pay digitally but manually — run the recurring-billing switch campaign.",
     "THIS QUARTER"),
    ("03", "Pay for word of mouth", "Referrals + returning tenants: 39% of tenants, 41% of "
     "$12.5M lifetime value, $14.6k per referred tenant. Launch referral credits & move-out win-back.",
     "REFERRAL CREDIT"),
    ("04", "Rebuild the website to book", "One web booking in 12 months while a third of "
     "tenants discover us online. Live availability, pricing, Reserve Now, local search.",
     "BOOKING FUNNEL"),
]
ay = 1.62
for n, t, d, tag in asks:
    rect(s, ML, ay, SW - 2 * ML, 1.16, fill=DARKCARD, line=DARKLINE, line_w=1, round_=0.06)
    rect(s, ML, ay + 0.06, 0.055, 1.04, fill=ACCENT)
    chip(s, ML + 0.24, ay + 0.36, n, dark=True)
    text(s, ML + 0.88, ay + 0.16, 3.15, 0.55, t, size=15.5, font=HEAD, bold=True,
         color=WHITE, line_spacing=1.0)
    text(s, ML + 4.15, ay + 0.14, 5.9, 0.9, d, size=10.5,
         color=RGBColor(0xC8, 0xC8, 0xC8), line_spacing=1.16, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, ML + 10.3, ay + 0.41, 1.88, 0.36, fill=RGBColor(0x2A, 0x12, 0x14),
         line=ACCENT_D, line_w=0.75, round_=0.5)
    text(s, ML + 10.3, ay + 0.41, 1.88, 0.36, tag, size=7.5, font=HEAD,
         bold=True, color=PINK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         spc=40, wrap=False)
    ay += 1.30
footer(s, 9, "Full detail in the Data Dashboard — occupancy, payments, acquisition, lead funnel & profile", dark=True)

out = "/home/user/storealldashboard/docs/storeall-data-story-jul-2026.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides.__iter__().__length_hint__() if hasattr(prs.slides,'__length_hint__') else 0) if False else len(prs.slides._sldIdLst))
